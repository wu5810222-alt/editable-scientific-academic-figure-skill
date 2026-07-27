#!/usr/bin/env python3
"""Audit whether a PPTX diagram is composed of native editable objects."""

import argparse
import json
import sys
from collections import Counter

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx")
    parser.add_argument("--max-pictures", type=int, default=0)
    parser.add_argument("--min-shapes", type=int, default=2)
    args = parser.parse_args()

    prs = Presentation(args.pptx)
    counts = Counter()
    text_objects = 0
    text_characters = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            counts[shape.shape_type.name] += 1
            text = getattr(shape, "text", "").strip()
            if text:
                text_objects += 1
                text_characters += len(text)

    pictures = counts[MSO_SHAPE_TYPE.PICTURE.name]
    total_shapes = sum(counts.values())
    failures = []
    if pictures > args.max_pictures:
        failures.append(f"picture count {pictures} exceeds allowed {args.max_pictures}")
    if total_shapes < args.min_shapes:
        failures.append(f"shape count {total_shapes} is below required {args.min_shapes}")

    report = {
        "file": args.pptx,
        "slides": len(prs.slides),
        "slide_size_inches": [prs.slide_width / 914400, prs.slide_height / 914400],
        "total_shapes": total_shapes,
        "shape_types": dict(sorted(counts.items())),
        "pictures": pictures,
        "text_objects": text_objects,
        "text_characters": text_characters,
        "object_level_editability_check": "pass" if not failures else "fail",
        "failures": failures,
    }
    print(json.dumps(report, indent=2, ensure_ascii=False))
    return 1 if failures else 0


if __name__ == "__main__":
    sys.exit(main())
